# Name & Study Nr.
# Deniz Isikli s215818
# Mark Nielsen s204434

import os
import pandas as pd
from tqdm import tqdm
from Util.util import Util
from pptx.util import Inches
from pptx import Presentation
import matplotlib.pyplot as plt
from sklearn.cluster import KMeans
from sklearn.decomposition import PCA
from sklearn.metrics import silhouette_score
from sklearn.preprocessing import StandardScaler

class VisualizationAndEvaluation():
    def __init__(self):
        """Run the visualization and evaluation of clustering results."""
        presentation = Presentation()

        steps = ["Loading & Scaling Data", "Plotting Elbow", "Calculating Silhouette", "Plotting KMeans PCA", "Plotting DBSCAN PCA"]
        progress = tqdm(total=len(steps), desc="Evaluation & Visualization", ncols=80)

        script_dir = os.path.dirname(os.path.abspath(__file__))
        file_path = os.path.join(script_dir, "Results", "TaskModelSklearn", "task_model_clusters.csv")
        os.makedirs(os.path.join(script_dir, "Results", "VisualizationAndEvaluation"), exist_ok=True)
        results_dir = os.path.join(script_dir, "Results", "VisualizationAndEvaluation")

        Util.remove_folder_content(results_dir)

        progress.set_postfix_str(steps[0])
        df = pd.read_csv(file_path)
        features = ['Recency', 'Frequency', 'Monetary']
        scaler = StandardScaler()
        X_scaled = scaler.fit_transform(df[features])
        progress.update(1)

        progress.set_postfix_str(steps[1])
        self.plot_elbow(X_scaled, presentation)
        progress.update(1)

        progress.set_postfix_str(steps[2])
        self.plot_silhouette(X_scaled, df['KMeansCluster'], presentation)
        progress.update(1)

        progress.set_postfix_str(steps[3])
        self.plot_pca(X_scaled, df['KMeansCluster'], "KMeans Clusters", presentation)
        progress.update(1)

        progress.set_postfix_str(steps[4])
        self.plot_pca(X_scaled, df['DBSCANCluster'], "DBSCAN Clusters", presentation)
        progress.update(1)

        presentation.save(os.path.join(results_dir, "Cluster_Evaluation_Presentation.pptx"))

        progress.close()
        print("Visualization and Evaluation completed.\n\n")

    def plot_elbow(self, X, presentation=None):
        """Plot the elbow method to determine the optimal number of clusters for KMeans.

        Args:
            X (ndarray): The input data.

        Outputs:
            - A plot showing the inertia (WCSS) for different numbers of clusters.
        """
        distortions = []
        for k in range(1, 10):
            km = KMeans(n_clusters=k, random_state=42)
            km.fit(X)
            distortions.append(km.inertia_)
        plt.plot(range(1, 10), distortions, marker='o')
        plt.xlabel('Number of Clusters')
        plt.ylabel('Inertia (WCSS)')
        plt.title('Elbow Method')
        filepath = "Results/VisualizationAndEvaluation/elbow_plot.png"
        plt.savefig(filepath)
        plt.close()
        if presentation:
            self.add_plot_to_slide(filepath, "Elbow Method for KMeans", presentation)

    def plot_silhouette(self, X, labels, presentation=None):
        """Plot the silhouette score for the clustering.

        Args:
            X (ndarray): The input data.
            labels (ndarray): The cluster labels.

        Outputs:
            - A plot showing the silhouette score for the clustering.
        """
        score = silhouette_score(X, labels)

        plt.bar(["Silhouette Score"], [score], color='skyblue')
        plt.ylim(0, 1)
        plt.title("Silhouette Score")
        plt.ylabel("Score")
        plt.xlabel("Clustering Method")
        plt.axhline(y=score, color='r', linestyle='--', label=f'Silhouette Score: {score:.3f}')
        plt.legend()
        filepath = "Results/VisualizationAndEvaluation/silhouette_score.png"
        plt.savefig(filepath)
        plt.close()
        if presentation:
            self.add_plot_to_slide(filepath, "Silhouette Score", presentation)

    def plot_pca(self, X, labels, title, presentation=None):
        """Plot the PCA components of the data.

        Args:
            X (ndarray): The input data.
            labels (ndarray): The cluster labels.
            title (str): The title for the plot.

        Outputs:
            - A scatter plot of the PCA components colored by cluster labels.
        """
        pca = PCA(n_components=2)
        components = pca.fit_transform(X)
        plt.scatter(components[:, 0], components[:, 1], c=labels, cmap='tab10', s=40)
        plt.title(title)
        plt.xlabel("PCA 1")
        plt.ylabel("PCA 2")
        filepath = f"Results/VisualizationAndEvaluation/{title.replace(' ', '_').lower()}.png"
        plt.savefig(filepath)
        plt.close()
        if presentation:
            self.add_plot_to_slide(filepath, title, presentation)

    def add_plot_to_slide(self, image_path, title, presentation):
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        slide.shapes.title.text = title
        slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), width=Inches(8))
